from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from utils import slugify, hex_to_rgb, is_dark_color, shorten_paragraph

def create_ppt_from_json(slides_data, theme_data, topic):
    filename = f"{slugify(topic)}.pptx"
    prs = Presentation()

    palette = theme_data.get("color_palette", ["#F1FAEE", "#A8DADC", "#457B9D"])
    font_name = theme_data.get("font_style", "Arial")

    background_hex = palette[0]
    if not is_dark_color(background_hex):
        background_hex = "#0F4C75"

    background_color = hex_to_rgb(background_hex)
    font_color = RGBColor(255, 255, 255) if is_dark_color(background_hex) else RGBColor(0, 0, 0)

    for i, slide_data in enumerate(slides_data):
        title = slide_data.get("title", "").strip()
        content = slide_data.get("content", [])

        if not title and not content:
            continue

        if title in ["Any Questions?", "Thank You"]:
            blank_slide = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide)

            left = top = Inches(2)
            width = height = Inches(6)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = title

            p = tf.paragraphs[0]
            p.font.size = Pt(50)
            p.font.italic = True
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p.font.color.rgb = font_color
            p.font.name = font_name

            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = background_color
            continue

        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title
        title_shape = slide.shapes.title
        if title_shape.has_text_frame:
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(32)
                    run.font.bold = True
                    run.font.name = font_name
                    run.font.color.rgb = font_color

        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        if layout == prs.slide_layouts[1]:
            tf = slide.placeholders[1].text_frame
            tf.clear()

            if title == "Agenda":
                for bullet in content:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(18)
                    p.font.name = font_name
                    p.font.color.rgb = font_color
            else:
                for bullet in content:
                    clean_text = bullet.replace("**", "").lstrip("• ").strip()
                    short_text = shorten_paragraph(clean_text)
                    if short_text:
                        p = tf.add_paragraph()
                        p.text = short_text
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.name = font_name
                        p.font.color.rgb = font_color
                        spacer = tf.add_paragraph()
                        spacer.text = ""

    prs.save(filename)
    return filename
